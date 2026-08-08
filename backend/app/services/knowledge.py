from __future__ import annotations
import csv, io, json, re
from pathlib import Path
from pypdf import PdfReader
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

ALLOWED = {'.txt', '.md', '.csv', '.json', '.pdf', '.docx', '.xlsx', '.pptx'}

def extract_text(filename: str, data: bytes) -> str:
    ext = Path(filename).suffix.lower()
    if ext not in ALLOWED:
        raise ValueError(f'Unsupported file type: {ext}')
    if ext in {'.txt', '.md'}:
        return data.decode('utf-8', errors='replace')
    if ext == '.json':
        obj = json.loads(data.decode('utf-8', errors='replace'))
        return json.dumps(obj, ensure_ascii=False, indent=2)
    if ext == '.csv':
        s = data.decode('utf-8-sig', errors='replace')
        rows = list(csv.reader(io.StringIO(s)))
        return '\n'.join(' | '.join(r) for r in rows)
    if ext == '.pdf':
        r = PdfReader(io.BytesIO(data))
        # No OCR fallback: a scanned/image-only PDF yields empty text here and the source will
        # simply have zero chunks rather than silently pretending to be indexed.
        return '\n'.join((p.extract_text() or '') for p in r.pages)
    if ext == '.docx':
        d = Document(io.BytesIO(data))
        return '\n'.join(p.text for p in d.paragraphs)
    if ext == '.xlsx':
        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        parts = []
        for ws in wb.worksheets:
            parts.append(f'# {ws.title}')
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    parts.append(' | '.join(cells))
        return '\n'.join(parts)
    if ext == '.pptx':
        prs = Presentation(io.BytesIO(data))
        parts = []
        for i, slide in enumerate(prs.slides, start=1):
            parts.append(f'# Slide {i}')
            for shape in slide.shapes:
                if getattr(shape, 'has_text_frame', False) and shape.text_frame.text.strip():
                    parts.append(shape.text_frame.text.strip())
        return '\n'.join(parts)
    return ''

def score_item(question: str, text: str) -> int:
    words = {w for w in re.findall(r'[\w\u0600-\u06ff]{3,}', question.lower())}
    low = text.lower()
    return sum(1 for w in words if w in low)

def top_context(question: str, items: list, limit: int = 8, max_chars: int = 50000):
    ranked = sorted(items, key=lambda x: score_item(question, x.content), reverse=True)
    chosen = [x for x in ranked if score_item(question, x.content) > 0][:limit] or ranked[:min(4, len(ranked))]
    chunks = []
    total = 0
    for x in chosen:
        text = x.content.strip()
        if total + len(text) > max_chars:
            text = text[:max(0, max_chars - total)]
        chunks.append({'id': x.id, 'title': x.title, 'source_type': x.source_type, 'text': text})
        total += len(text)
        if total >= max_chars:
            break
    return chunks
